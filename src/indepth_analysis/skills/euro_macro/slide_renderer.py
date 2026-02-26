"""PowerPoint dashboard slide renderer for Euro Macro reports."""

import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from indepth_analysis.models.euro_macro import EuroMacroReport, ReportSection

logger = logging.getLogger(__name__)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
METRIC_BLUE = RGBColor(0x0D, 0x47, 0xA1)

# Layout constants
MARGIN = Inches(0.3)
TITLE_HEIGHT = Inches(0.7)
METRIC_ROW_TOP = TITLE_HEIGHT + Inches(0.15)
METRIC_ROW_HEIGHT = Inches(0.85)
GRID_TOP = METRIC_ROW_TOP + METRIC_ROW_HEIGHT + Inches(0.15)
FOOTER_HEIGHT = Inches(0.4)
GRID_HEIGHT = SLIDE_HEIGHT - GRID_TOP - FOOTER_HEIGHT - Inches(0.1)

# Section heading short labels
SECTION_SHORT_LABELS = {
    "통화정책": "통화정책",
    "경제성장": "경제성장",
    "물가": "물가",
    "금융시장": "금융시장",
    "정치": "정치현황",
    "무역": "무역·관세",
    "리스크": "리스크",
    "전망": "전망·시사점",
}


def _short_label(heading: str) -> str:
    """Extract a short label from a section heading."""
    for key, label in SECTION_SHORT_LABELS.items():
        if key in heading:
            return label
    # Strip leading number prefix like "1. "
    return re.sub(r"^\d+\.\s*", "", heading)[:8]


# Metric extraction patterns
_METRIC_PATTERNS: list[tuple[str, str, str]] = [
    ("ECB 금리", r"예금금리\s*(\d+\.?\d*%)", "1. 통화정책"),
    ("GDP 성장률", r"(?:실질\s*)?GDP\s*성장률[을를]?\s*(\d+\.?\d*%)", "2. 경제성장"),
    ("CPI", r"(?:연간\s*)?인플레이션[율률]?[은는이가]?\s*(\d+\.?\d*%)", "3. 물가"),
    ("PMI", r"제조업\s*PMI[^0-9]*(\d+\.?\d*)", "2. 경제성장"),
    ("실업률", r"실업률[은는이가]?\s*(\d+\.?\d*%)", "2. 경제성장"),
    ("EUR/USD", r"EUR/?USD[^0-9~]*(~?\d+\.?\d*)", "4. 금융시장"),
]


def extract_key_metrics(report: EuroMacroReport) -> list[dict[str, str]]:
    """Extract key metrics from report section content.

    Returns a list of {"label": str, "value": str} dicts.
    """
    # Build a lookup of section content by number prefix
    section_content: dict[str, str] = {}
    full_text = ""
    for section in report.sections:
        full_text += section.content + "\n"
        for prefix in ["1.", "2.", "3.", "4.", "5.", "6.", "7.", "8."]:
            if section.heading.startswith(prefix) or prefix[0] in section.heading[:3]:
                section_content[prefix] = section.content

    metrics: list[dict[str, str]] = []
    for label, pattern, section_hint in _METRIC_PATTERNS:
        # Search preferred section first, then full text
        hint_prefix = section_hint.split(".")[0] + "."
        search_text = section_content.get(hint_prefix, "") + "\n" + full_text
        match = re.search(pattern, search_text)
        value = match.group(1) if match else "N/A"
        metrics.append({"label": label, "value": value})

    return metrics


def extract_section_bullets(section: ReportSection, max_bullets: int = 3) -> list[str]:
    """Extract first N sentences from section content as bullet points.

    Each bullet is truncated to ~80 characters.
    """
    # Split on sentence-ending punctuation (Korean period, regular period)
    text = section.content.strip()
    # Remove source citations for cleaner bullets
    text = re.sub(r"\s*\[출처:[^\]]*\]", "", text)

    sentences = re.split(r"(?<=[.다])\s+", text)
    bullets: list[str] = []
    for s in sentences:
        s = s.strip()
        if not s:
            continue
        if len(s) > 80:
            s = s[:77] + "..."
        bullets.append(s)
        if len(bullets) >= max_bullets:
            break

    return bullets


def _add_shape(slide, left, top, width, height, fill_color=None):
    """Add a rectangle shape to the slide."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()  # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _set_text(
    shape,
    text,
    font_size=Pt(10),
    bold=False,
    color=DARK_GRAY,
    alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(72000)
    tf.margin_right = Emu(72000)
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

    shape.text_frame.paragraphs[0].space_after = Pt(0)
    shape.text_frame.paragraphs[0].space_before = Pt(0)

    # Vertical anchor
    tf.word_wrap = True


def render_slide(report: EuroMacroReport) -> Presentation:
    """Create a single dashboard slide from the report.

    Returns a python-pptx Presentation object with one slide.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    content_width = SLIDE_WIDTH - 2 * MARGIN

    # --- Title bar ---
    title_shape = _add_shape(
        slide, MARGIN, Inches(0.1), content_width, TITLE_HEIGHT, DARK_BLUE
    )
    _set_text(
        title_shape,
        report.title,
        font_size=Pt(22),
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- Key metrics row ---
    metrics = extract_key_metrics(report)
    num_metrics = len(metrics)
    metric_gap = Inches(0.12)
    total_gap = metric_gap * (num_metrics - 1)
    metric_width = int((content_width - total_gap) / num_metrics)

    for i, m in enumerate(metrics):
        left = MARGIN + i * (metric_width + metric_gap)
        box = _add_shape(
            slide, left, METRIC_ROW_TOP, metric_width, METRIC_ROW_HEIGHT, LIGHT_BLUE
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(54000)
        tf.margin_right = Emu(54000)
        tf.margin_top = Emu(36000)
        tf.margin_bottom = Emu(18000)

        # Label paragraph
        p_label = tf.paragraphs[0]
        p_label.text = m["label"]
        p_label.font.size = Pt(10)
        p_label.font.bold = True
        p_label.font.color.rgb = MID_GRAY
        p_label.alignment = PP_ALIGN.CENTER
        p_label.space_after = Pt(2)

        # Value paragraph
        p_value = tf.add_paragraph()
        p_value.text = m["value"]
        p_value.font.size = Pt(18)
        p_value.font.bold = True
        p_value.font.color.rgb = METRIC_BLUE
        p_value.alignment = PP_ALIGN.CENTER
        p_value.space_before = Pt(0)

    # --- Section grid (2 rows × 4 columns) ---
    sections = report.sections[:8]
    cols = 4
    rows = 2
    grid_gap = Inches(0.1)
    cell_width = int((content_width - grid_gap * (cols - 1)) / cols)
    cell_height = int((GRID_HEIGHT - grid_gap * (rows - 1)) / rows)

    for idx, section in enumerate(sections):
        row = idx // cols
        col = idx % cols
        left = MARGIN + col * (cell_width + grid_gap)
        top = GRID_TOP + row * (cell_height + grid_gap)

        box = _add_shape(slide, left, top, cell_width, cell_height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(72000)
        tf.margin_right = Emu(72000)
        tf.margin_top = Emu(36000)
        tf.margin_bottom = Emu(36000)

        # Section heading
        label = _short_label(section.heading)
        p_heading = tf.paragraphs[0]
        p_heading.text = label
        p_heading.font.size = Pt(11)
        p_heading.font.bold = True
        p_heading.font.color.rgb = DARK_BLUE
        p_heading.space_after = Pt(4)

        # Bullets
        bullets = extract_section_bullets(section)
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(7.5)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(1)
            p.space_after = Pt(1)

    # --- Footer ---
    footer_top = SLIDE_HEIGHT - FOOTER_HEIGHT - Inches(0.05)
    footer_shape = _add_shape(slide, MARGIN, footer_top, content_width, FOOTER_HEIGHT)
    _set_text(
        footer_shape,
        "출처: KCIF, ECB, Eurostat, S&P Global, IMF, OECD",
        font_size=Pt(8),
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    return prs


def save_slide(report: EuroMacroReport, output_dir: str = "reports") -> Path:
    """Render and save the dashboard slide.

    Returns the path to the saved .pptx file.
    """
    prs = render_slide(report)

    dir_path = Path(output_dir) / "euro_macro"
    dir_path.mkdir(parents=True, exist_ok=True)

    filename = f"{report.year}-{report.month:02d}.pptx"
    filepath = dir_path / filename
    prs.save(str(filepath))

    logger.info("Slide saved to %s", filepath)
    return filepath
