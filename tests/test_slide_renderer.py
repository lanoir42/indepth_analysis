"""Tests for the PowerPoint slide renderer."""

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from indepth_analysis.models.euro_macro import (
    AgentResult,
    EuroMacroReport,
    ReportSection,
)
from indepth_analysis.skills.euro_macro.slide_renderer import (
    extract_key_metrics,
    extract_section_bullets,
    render_slide,
    save_slide,
)


def _make_report() -> EuroMacroReport:
    """Build a minimal EuroMacroReport with realistic content."""
    return EuroMacroReport(
        year=2026,
        month=2,
        title="2026년 2월 월간 유럽 거시경제 현황",
        sections=[
            ReportSection(
                heading="1. 통화정책 및 금리 동향",
                content=(
                    "ECB는 예금금리 2.0%를 유지했다. "
                    "라가르드 총재는 금리가 적절한 수준이라고 강조했다. "
                    "시장은 연내 동결을 예상한다."
                ),
            ),
            ReportSection(
                heading="2. 경제성장 및 고용",
                content=(
                    "유로존 실질 GDP 성장률을 1.2%로 전망한다. "
                    "제조업 PMI는 50.8로 확장 국면에 진입했다. "
                    "실업률은 6.3%를 기록했다."
                ),
            ),
            ReportSection(
                heading="3. 물가 및 인플레이션",
                content=(
                    "1월 연간 인플레이션율은 1.7%로 하락했다. "
                    "근원 인플레이션은 2.2%를 기록했다. "
                    "에너지 가격은 하락세를 보였다."
                ),
            ),
            ReportSection(
                heading="4. 금융시장 동향",
                content=(
                    "EUR/USD 환율은 ~1.04 수준이다. "
                    "Eurostoxx50은 혼조세를 보였다. "
                    "장기금리는 안정적이었다."
                ),
            ),
            ReportSection(
                heading="5. EU 및 주요국 정치 현황",
                content="독일 CDU/CSU-SPD 연정이 개혁을 추진 중이다.",
            ),
            ReportSection(
                heading="6. 무역 및 관세 현황",
                content="트럼프 대통령은 15% 관세를 부과했다.",
            ),
            ReportSection(
                heading="7. 주요 리스크 요인",
                content="AI 버블 우려와 위험자산 가격 조정이 주요 리스크다.",
            ),
            ReportSection(
                heading="8. 향후 전망 및 시사점",
                content="유로존 경제는 완만한 회복세를 이어갈 전망이다.",
            ),
        ],
        agent_results=[
            AgentResult(agent_name="KCIF", findings=[], search_queries=[]),
        ],
        model_used="test-model",
        total_findings=0,
    )


class TestExtractKeyMetrics:
    def test_extracts_all_six_metrics(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        assert len(metrics) == 6
        labels = [m["label"] for m in metrics]
        assert "ECB 금리" in labels
        assert "GDP 성장률" in labels
        assert "CPI" in labels
        assert "PMI" in labels
        assert "실업률" in labels
        assert "EUR/USD" in labels

    def test_ecb_rate_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        ecb = next(m for m in metrics if m["label"] == "ECB 금리")
        assert ecb["value"] == "2.0%"

    def test_gdp_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        gdp = next(m for m in metrics if m["label"] == "GDP 성장률")
        assert gdp["value"] == "1.2%"

    def test_cpi_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        cpi = next(m for m in metrics if m["label"] == "CPI")
        assert cpi["value"] == "1.7%"

    def test_pmi_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        pmi = next(m for m in metrics if m["label"] == "PMI")
        assert pmi["value"] == "50.8"

    def test_unemployment_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        unemp = next(m for m in metrics if m["label"] == "실업률")
        assert unemp["value"] == "6.3%"

    def test_eurusd_value(self):
        report = _make_report()
        metrics = extract_key_metrics(report)
        fx = next(m for m in metrics if m["label"] == "EUR/USD")
        assert fx["value"] == "~1.04"

    def test_missing_metric_returns_na(self):
        report = EuroMacroReport(
            year=2026,
            month=2,
            title="Empty Report",
            sections=[
                ReportSection(heading="1. 통화정책", content="내용 없음"),
            ],
        )
        metrics = extract_key_metrics(report)
        na_count = sum(1 for m in metrics if m["value"] == "N/A")
        assert na_count >= 1


class TestExtractSectionBullets:
    def test_returns_max_bullets(self):
        section = ReportSection(
            heading="1. 통화정책",
            content=(
                "첫 번째 문장이다. "
                "두 번째 문장이다. "
                "세 번째 문장이다. "
                "네 번째 문장이다."
            ),
        )
        bullets = extract_section_bullets(section, max_bullets=3)
        assert len(bullets) == 3

    def test_truncates_long_bullets(self):
        long_text = "이것은 매우 긴 문장입니다 " * 10 + "끝이다."
        section = ReportSection(heading="test", content=long_text)
        bullets = extract_section_bullets(section, max_bullets=1)
        assert len(bullets) == 1
        assert len(bullets[0]) <= 83  # 80 + "..."

    def test_strips_source_citations(self):
        section = ReportSection(
            heading="test",
            content="ECB는 금리를 동결했다 [출처: ECB]. 시장은 안정적이다.",
        )
        bullets = extract_section_bullets(section)
        for b in bullets:
            assert "[출처:" not in b

    def test_empty_content(self):
        section = ReportSection(heading="test", content="")
        bullets = extract_section_bullets(section)
        assert bullets == []


class TestRenderSlide:
    def test_produces_valid_presentation(self):
        report = _make_report()
        prs = render_slide(report)
        assert isinstance(prs, PresentationType)

    def test_has_one_slide(self):
        report = _make_report()
        prs = render_slide(report)
        assert len(prs.slides) == 1

    def test_slide_is_widescreen(self):
        report = _make_report()
        prs = render_slide(report)
        # 16:9 ratio check
        ratio = prs.slide_width / prs.slide_height
        assert abs(ratio - 16 / 9) < 0.01


class TestSaveSlide:
    def test_writes_file(self, tmp_path):
        report = _make_report()
        filepath = save_slide(report, output_dir=str(tmp_path))
        assert filepath.exists()
        assert filepath.suffix == ".pptx"
        assert filepath.name == "2026-02.pptx"

    def test_file_is_valid_pptx(self, tmp_path):
        report = _make_report()
        filepath = save_slide(report, output_dir=str(tmp_path))
        prs = Presentation(str(filepath))
        assert len(prs.slides) == 1

    def test_creates_directory(self, tmp_path):
        report = _make_report()
        output_dir = tmp_path / "nested" / "output"
        filepath = save_slide(report, output_dir=str(output_dir))
        assert filepath.exists()
        assert filepath.parent == output_dir / "euro_macro"
